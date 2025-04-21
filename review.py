import os
import re
import json
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from src.medrag import MedRAG

client = OpenAI(
    api_key=os.environ.get("OPENROUTER_API_KEY"),
    base_url="https://openrouter.ai/api/v1"
)
openai_client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
)
MODEL = "deepseek/deepseek-r1"

# client = OpenAI(
#     api_key=os.environ.get("OPENAI_API_KEY"),
# )
# MODEL = "gpt-4o"

def sanitize_filename(s):
    """Return a filename-safe version of the subject string."""
    s = s.lower()
    s = re.sub(r"[,\(\)]", "", s)  # remove commas and parentheses
    s = re.sub(r"\s+", "_", s)      # replace whitespace with underscore
    return s

def get_high_yield_concepts(subject, num_questions):
    """
    Ask the LLM to generate a numbered list of high-yield USMLE Step 1 concepts for the given subject area.
    The prompt covers diseases, pathophysiology, pharmacology, physiology, anatomy, microbiology, embryology, etc.
    """
    prompt = (
        f"Please generate a list of {num_questions} high-yield must-know USMLE Step 1 {subject} concepts. "
        "Include topics that cover diseases, pathophysiology, pharmacology, physiology, anatomy, microbiology, embryology, etc as appropriate. "
        "Return the list in a numbered format (one concept per line)."
    )
    response = openai_client.chat.completions.create(
        model='gpt-4o',
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    concepts_text = response.choices[0].message.content
    # Parse the numbered list
    concepts = []
    for line in concepts_text.splitlines():
        line = line.strip()
        if line:
            # Remove numbering if present
            parts = line.split('.', 1)
            if len(parts) > 1:
                concept = parts[1].strip()
            else:
                concept = line
            concepts.append(concept)
    return concepts[:num_questions]

def add_slide_with_content(prs, question, highlight=False, include_explanation=False):
    """
    Create a slide with the content for a single MCQ.
    - If highlight is True, the correct answer (matching question['correct_answer']) will be highlighted.
    - If include_explanation is True, an additional paragraph with the explanation is appended.
    """
    slide_layout = prs.slide_layouts[6]  # use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    # Create a textbox covering most of the slide
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4))
    tf = textbox.text_frame
    # Clear any default text (if present)
    tf.clear()
    tf.word_wrap = True

    # Question text
    p = tf.paragraphs[0]
    p.text = question["question"] + '\n'
    p.font.size = Pt(16)
    p.font.name = "Proxima Nova Regular"
    p.font.color.rgb = RGBColor(54, 54, 54)
    p.line_spacing = Pt(16*1.15)  # 18pt line height

    # Answer choices
    for choice in question["options"].keys():
        p = tf.add_paragraph()
        p.font.size = Pt(16)
        p.font.name = "Proxima Nova Regular"
        p.line_spacing = Pt(16*1.15)  # 18pt line height
        p.font.color.rgb = RGBColor(54, 54, 54)
        p.text = f'{choice}: {question["options"][choice]}'
        if highlight and choice.strip() == question["correct_answer"].strip():
            p.font.color.rgb = RGBColor(255, 0, 0)  # Dark red for highlighting

    p = tf.add_paragraph()
    p.text = ""

    # If explanation is requested, add it below the answer choices.
    if include_explanation:
        p = tf.add_paragraph()
        p.text = f"{question["correct_answer"]}: {question['explanation'][question["correct_answer"]].replace('Correct. ', '')}\n"
        p.font.color.rgb = RGBColor(255, 0, 0)  # Dark red for highlighting
        p.font.size = Pt(12)
        p.font.name = "Proxima Nova Regular"
        p.line_spacing = Pt(12*1.15)  # 18pt line height

        for choice in question["options"].keys():
            if choice != question["correct_answer"]:
                p = tf.add_paragraph()
                p.text = f"{choice}: {question['explanation'][choice].replace('Incorrect. ', '')}\n"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(54, 54, 54)
                p.font.name = "Proxima Nova Regular"
                p.line_spacing = Pt(12*1.15)  # 18pt line height

    return slide

def create_ppt_for_subject(subject, questions):
    """
    Create a PowerPoint presentation for the given subject group.
    For each MCQ, create three slides:
      1. Slide 1: The vignette and answer choices (no highlighting).
      2. Slide 2: Duplicate of slide 1 with the correct answer highlighted.
      3. Slide 3: Duplicate of slide 2 with the explanation appended.
    """
    prs = Presentation()

    for idx, q in enumerate(questions, start=1):
        if 'json' in q:
            q = json.loads(q.split('json')[1].split('```')[0])
        else:
            q = json.loads(q)
        # Slide 1: Basic content
        add_slide_with_content(prs, q, idx, highlight=False, include_explanation=False)
        # Slide 2: Highlight the correct answer
        add_slide_with_content(prs, q, idx, highlight=True, include_explanation=True)
    
    filename = f"{sanitize_filename(subject)}.pptx"
    prs.save(os.path.join('deepseek', filename))
    print(f"Presentation for '{subject}' created successfully as '{filename}'.")

def main():
    # Define subject groups and the proportional number of questions for each
    subject_groups = {
        "biochemistry, genetics, pharmacology, poisoning and environmental exposure": 50,
        "allergy, immunology, infectious disease, microbiology": 100,
        "cardiology": 100,
        "hematology, oncology": 50,
        "respiratory, pulmonology": 50,
        "renal, urology": 50,
        "endocrine": 50,
        "reproductive (male and female, pregnancy)": 50,
        "gastrointestinal": 50,
        "dermatology": 50,
        "orthopedics, rheumatology": 50,
        "neurology, ophthalmology and ent": 100,
        "psychiatry": 50,
        "biostatistics": 20
    }

    medrag = MedRAG(llm_name="deepseek/deepseek-r1", rag=True, retriever_name="MedCPT", corpus_name="Textbooks")
    
    # Loop over each subject group and generate the corresponding content and PPT.
    for subject, num_questions in subject_groups.items():
        print(f"\n=== Generating content for subject: {subject} ({num_questions} questions) ===")
        # Step 1: Generate high-yield concepts for the subject area
        concepts = get_high_yield_concepts(subject, num_questions)
        print(f"Retrieved {len(concepts)} concepts for {subject}.")
        
        subject_questions = []
        # Step 2 & 3: For each concept, generate a clinical vignette and then the corresponding MCQ
        for idx, concept in enumerate(concepts, start=1):
            print(f"Generating content for concept {idx}: {concept}")
            mcq = medrag.generate_usmle_question(concept)
            if mcq is not None:
                subject_questions.append(mcq)
        
        # Step 4: Create and save the PowerPoint presentation for this subject
        create_ppt_for_subject(subject, subject_questions)

if __name__ == "__main__":
    main()
